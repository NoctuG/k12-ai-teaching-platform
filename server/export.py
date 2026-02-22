#!/usr/bin/env python3
import sys
import json
import re
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt


def parse_markdown(markdown: str):
    blocks = []
    current_list = []
    for raw in markdown.splitlines():
        line = raw.strip()
        if not line:
            if current_list:
                blocks.append({"type": "list", "items": current_list[:]})
                current_list = []
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
        list_match = re.match(r"^[-*]\s+(.+)$", line)

        if heading_match:
            if current_list:
                blocks.append({"type": "list", "items": current_list[:]})
                current_list = []
            blocks.append({"type": "heading", "level": len(heading_match.group(1)), "text": heading_match.group(2).strip()})
        elif list_match:
            current_list.append(list_match.group(1).strip())
        else:
            if current_list:
                blocks.append({"type": "list", "items": current_list[:]})
                current_list = []
            blocks.append({"type": "paragraph", "text": line})

    if current_list:
        blocks.append({"type": "list", "items": current_list[:]})
    return blocks


def export_to_docx(markdown: str, title: str, output_path: str):
    doc = Document()

    title_paragraph = doc.add_paragraph()
    title_run = title_paragraph.add_run(title)
    title_run.font.size = Pt(18)
    title_run.bold = True
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    for block in parse_markdown(markdown):
        if block["type"] == "heading":
            level = min(block["level"], 4)
            p = doc.add_heading(block["text"], level=level)
            p.paragraph_format.space_after = Pt(6)
        elif block["type"] == "list":
            for item in block["items"]:
                p = doc.add_paragraph(item, style="List Bullet")
                p.paragraph_format.space_after = Pt(2)
        else:
            p = doc.add_paragraph(block["text"])
            p.paragraph_format.line_spacing = 1.4
            p.paragraph_format.space_after = Pt(8)

    doc.save(output_path)


def chunk(items, n):
    return [items[i:i + n] for i in range(0, len(items), n)]


def export_to_pptx(markdown: str, title: str, output_path: str, resource_type: str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    blocks = parse_markdown(markdown)

    # courseware strategy: heading => slide, section body paginated into multiple slides
    if resource_type == "courseware":
        sections = []
        current = {"title": "课程内容", "bullets": []}
        for block in blocks:
            if block["type"] == "heading" and block["level"] <= 2:
                if current["bullets"]:
                    sections.append(current)
                current = {"title": block["text"], "bullets": []}
            elif block["type"] == "list":
                current["bullets"].extend(block["items"])
            else:
                current["bullets"].append(block.get("text", ""))

        if current["bullets"]:
            sections.append(current)

        for section in sections:
            pages = chunk([x for x in section["bullets"] if x], 6)
            for idx, page in enumerate(pages):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                suffix = f"（第{idx+1}页）" if len(pages) > 1 else ""
                slide.shapes.title.text = f"{section['title']}{suffix}"
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for line in page:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
                    p.font.size = PptPt(20)
    else:
        bullets = []
        for block in blocks:
            if block["type"] == "list":
                bullets.extend(block["items"])
            else:
                bullets.append(block.get("text", ""))
        for idx, page in enumerate(chunk([x for x in bullets if x], 8), start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"内容 {idx}"
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for line in page:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = PptPt(18)

    prs.save(output_path)


def main():
    if len(sys.argv) < 5:
        print(json.dumps({"error": "Usage: export.py <format> <title> <markdown> <output_path> [resource_type]"}))
        sys.exit(1)

    format_type = sys.argv[1]
    title = sys.argv[2]
    markdown = sys.argv[3]
    output_path = sys.argv[4]
    resource_type = sys.argv[5] if len(sys.argv) > 5 else ""

    try:
        if format_type == "docx":
            export_to_docx(markdown, title, output_path)
        elif format_type == "pptx":
            export_to_pptx(markdown, title, output_path, resource_type)
        else:
            raise ValueError(f"Unsupported format: {format_type}")

        print(json.dumps({"success": True, "path": output_path}))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
